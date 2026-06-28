"""Build the EHL Paris Medical Image Retrieval presentation (dark theme, PPTX).

Research-piece structure (~5 min):
  1. Title (no numbers)
  2. The data  -> why contrastive is the natural method
  3. Architecture: cross-modal encoding + classical realization + Hungarian
  4. Data leakage: the discovery, how we removed it, honest boundaries
  5. What did NOT work: augmentation
  6. Metrics: honest leak-free results
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette ----
BG     = RGBColor(0x0D, 0x11, 0x17)
PANEL  = RGBColor(0x16, 0x1B, 0x22)
PANEL2 = RGBColor(0x1C, 0x23, 0x2D)
WHITE  = RGBColor(0xE6, 0xED, 0xF3)
MUTE   = RGBColor(0x9D, 0xA7, 0xB3)
CYAN   = RGBColor(0x22, 0xD3, 0xEE)
MAG    = RGBColor(0xF4, 0x72, 0xB6)
LIME   = RGBColor(0xA3, 0xE6, 0x35)
AMBER  = RGBColor(0xFB, 0xBF, 0x24)
RED    = RGBColor(0xF8, 0x71, 0x71)
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)

FONT = "Segoe UI"
MONO = "Consolas"
EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    r._element.addprevious(r._element)  # keep at back
    return s


def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, round_=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (txt,size,color,bold,font)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, *f) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
            r.font.name = f[0] if f else FONT
    return tb


def header(s, kicker, kcolor, title_txt):
    box(s, 0.0, 0.0, 0.16, 7.5, fill=kcolor)
    text(s, 0.6, 0.42, 12.2, 0.4, [[(kicker.upper(), 13, kcolor, True)]])
    text(s, 0.6, 0.74, 12.2, 1.0, [[(title_txt, 30, WHITE, True)]])
    box(s, 0.62, 1.5, 1.1, 0.05, fill=kcolor)


def bullets(s, x, y, w, h, items, size=15, gap=10):
    runs = []
    for it in items:
        if isinstance(it, tuple):
            mark, txt, col = it
        else:
            mark, txt, col = "›", it, CYAN
        runs.append([(mark + "  ", size, col, True), (txt, size, WHITE, False)])
    text(s, x, y, w, h, runs, space_after=gap, line_spacing=1.08)


def arrow(s, x, y, w=0.36, col=MUTE):
    text(s, x, y, w, 0.5, [[("→", 22, col, True)]], align=PP_ALIGN.CENTER)


# ============================================================ 1. TITLE (no numbers)
s = slide()
box(s, 0, 0, 13.333, 0.16, fill=CYAN)
box(s, 0, 7.34, 13.333, 0.16, fill=MAG)
text(s, 0.9, 2.35, 11.5, 1.4,
     [[("Cross-Modal Brain MRI Retrieval", 46, WHITE, True)]])
text(s, 0.9, 3.45, 11.5, 0.7,
     [[("Matching T1 post-contrast queries to their T2 targets across three datasets", 19, CYAN, False)]])
text(s, 0.92, 4.55, 11.5, 0.5,
     [[("EHL Paris  ·  Medical Image Retrieval Challenge", 15, MUTE, True)]])
text(s, 0.92, 5.5, 11.5, 0.5,
     [[("A research walkthrough:  the data → our method → a data-leak we uncovered → honest metrics", 15, MUTE, False)]])

# ============================================================ 2. THE DATA -> WHY CONTRASTIVE
s = slide()
header(s, "The data", CYAN, "Why this problem is built for contrast")
text(s, 0.65, 1.9, 12.1, 1.0, [[
    ("Each patient gives one T1c (query) and one T2w (target) scan. The mapping is a clean "
     "bijection: every query has exactly one true target, and ", 16, WHITE, False),
    ("every other pair is a genuine non-match.", 16, CYAN, True)]],
    line_spacing=1.2)

# left: the bijection picture
box(s, 0.65, 3.1, 5.9, 3.4, fill=PANEL, round_=True)
text(s, 0.95, 3.3, 5.4, 0.4, [[("ONE-TO-ONE STRUCTURE", 12, CYAN, True)]])
for i in range(3):
    yy = 3.95 + i * 0.72
    box(s, 1.1, yy, 1.6, 0.5, fill=PANEL2, round_=True)
    text(s, 1.1, yy, 1.6, 0.5, [[("T1c q%d" % (i + 1), 12, CYAN, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    arrow(s, 2.85, yy + 0.02, col=LIME)
    box(s, 3.55, yy, 1.6, 0.5, fill=PANEL2, round_=True)
    text(s, 3.55, yy, 1.6, 0.5, [[("T2w t%d" % (i + 1), 12, MAG, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 1.1, 6.15, 5.0, 0.4, [[("exactly one true link per row — all off-diagonal pairs are negatives", 11, MUTE, False)]])

# right: why that is ideal for contrast
box(s, 6.85, 3.1, 5.9, 3.4, fill=PANEL, round_=True)
text(s, 7.15, 3.3, 5.4, 0.4, [[("IDEAL CONTRASTIVE SETUP", 12, LIME, True)]])
bullets(s, 7.15, 3.85, 5.4, 2.6, [
    ("✓", "Positives are given for free: the labelled query–target pairs.", LIME),
    ("✓", "Negatives are clean: any non-paired (q, t) is truly unrelated — no label noise.", LIME),
    ("✓", "Cross-modal: pull T1c and its T2w together, push everything else apart.", LIME),
    ("✓", "The bijection also lets us solve matching globally (Hungarian) at the end.", CYAN),
], size=12.5, gap=10)

# ============================================================ 3. ARCHITECTURE
s = slide()
header(s, "Method", VIOLET, "Cross-modal encoding → assignment")
# two-encoder contrastive schematic
text(s, 0.65, 1.85, 12.1, 0.4, [[("The architecture: two encoders, one shared embedding space", 14, VIOLET, True)]])
box(s, 0.9, 2.4, 2.3, 0.9, fill=PANEL2, round_=True); box(s, 0.9, 2.4, 2.3, 0.08, fill=CYAN)
text(s, 0.9, 2.4, 2.3, 0.9, [[("T1c encoder", 13, CYAN, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, 0.9, 3.55, 2.3, 0.9, fill=PANEL2, round_=True); box(s, 0.9, 3.55, 2.3, 0.08, fill=MAG)
text(s, 0.9, 3.55, 2.3, 0.9, [[("T2w encoder", 13, MAG, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
arrow(s, 3.35, 2.95, col=MUTE); arrow(s, 3.35, 3.6, col=MUTE)
box(s, 4.05, 2.7, 2.6, 1.45, fill=PANEL, round_=True, line=VIOLET, line_w=1.5)
text(s, 4.05, 2.7, 2.6, 1.45, [[("shared", 14, VIOLET, True)], [("embedding space", 12, WHITE, False)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
arrow(s, 6.8, 3.25, col=LIME)
box(s, 7.5, 2.7, 2.5, 1.45, fill=PANEL, round_=True); box(s, 7.5, 2.7, 2.5, 0.08, fill=AMBER)
text(s, 7.5, 2.7, 2.5, 1.45, [[("cosine", 14, AMBER, True)], [("similarity matrix", 12, WHITE, False)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
arrow(s, 10.15, 3.25, col=LIME)
box(s, 10.85, 2.7, 1.9, 1.45, fill=PANEL, round_=True); box(s, 10.85, 2.7, 1.9, 0.08, fill=LIME)
text(s, 10.85, 2.7, 1.9, 1.45, [[("Hungarian", 13, LIME, True)], [("one-to-one", 11, WHITE, False)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# realization details
box(s, 0.65, 4.55, 5.95, 1.95, fill=PANEL, round_=True)
text(s, 0.95, 4.72, 5.4, 0.4, [[("HOW WE REALIZE IT", 12, CYAN, True)]])
bullets(s, 0.95, 5.15, 5.5, 1.3, [
    ("›", "Neighbor-aware voxel feature: each voxel carries its local neighborhood, not just its own intensity.", CYAN),
    ("›", "Volumes registered into one common template frame so the cross-modal map transfers across subjects.", CYAN),
], size=11.5, gap=8)
box(s, 6.8, 4.55, 5.95, 1.95, fill=PANEL, round_=True)
text(s, 7.1, 4.72, 5.4, 0.4, [[("HUNGARIAN ASSIGNMENT", 12, LIME, True)]])
bullets(s, 7.1, 5.15, 5.5, 1.3, [
    ("★", "Exploit the bijection: solve the optimal one-to-one matching, force each assigned target to rank 1.", LIME),
    ("★", "Measured ablation: 0.651 → 0.718 MRR (+0.068) — the cheapest score lever in the pipeline.", LIME),
], size=11.5, gap=8)

# ============================================================ 4. DATA LEAKAGE
s = slide()
header(s, "What we uncovered", RED, "The leaderboard has a data leak")
text(s, 0.65, 1.9, 12.1, 1.0, [[
    ("The near-1.0 scores at the top are not a model — they are a leak. ", 16, RED, True),
    ("The competition scans come from the public BraTS-GLI library: each image can be re-identified "
     "by subject, which recovers its true pair.", 16, WHITE, False)]],
    line_spacing=1.2)
box(s, 0.65, 3.1, 5.9, 3.4, fill=PANEL, round_=True)
text(s, 0.95, 3.3, 5.4, 0.4, [[("WHAT WE DID", 12, AMBER, True)]])
bullets(s, 0.95, 3.8, 5.4, 2.5, [
    ("›", "Identified query→T1c and gallery→T2w against public BraTS; two independent methods agreed 19/19.", AMBER),
    ("›", "Reproduced the leak end-to-end: public score 0.931 — proof the board is gameable.", AMBER),
    ("›", "Then removed it and report honest, leak-free numbers.", LIME),
], size=12, gap=9)
box(s, 6.85, 3.1, 5.9, 3.4, fill=PANEL, round_=True)
text(s, 7.15, 3.3, 5.4, 0.4, [[("HONEST BOUNDARIES", 12, CYAN, True)]])
bullets(s, 7.15, 3.8, 5.4, 2.5, [
    ("→", "d3 = 1.0 is an affine / FOV-shape leak (the target is resliced into the query's frame).", CYAN),
    ("→", "d1 barely moves — a shared grid means no shape fingerprint (already 0.964).", CYAN),
    ("→", "d2 keeps an elastic warp that rigid registration cannot fully undo — the real wall.", CYAN),
], size=12, gap=9)

# ============================================================ 5. AUGMENTATION (did NOT work)
s = slide()
header(s, "What did NOT work", AMBER, "Augmentation to train a learned model")
text(s, 0.65, 1.9, 12.1, 0.9, [[
    ("Training data lives in a single domain, so we tried geometric + contrast augmentation to train a "
     "learned contrastive encoder. ", 15, WHITE, False),
    ("Every variant made the score worse.", 15, RED, True)]],
    line_spacing=1.2)
box(s, 0.65, 3.0, 5.9, 3.3, fill=PANEL, round_=True)
text(s, 0.95, 3.2, 5.4, 0.4, [[("WHAT WE TRIED", 12, AMBER, True)]])
bullets(s, 0.95, 3.7, 5.4, 2.5, [
    ("•", "Contrastive 3D model on augmented pairs → holdout MRR ≈ 0.04.", RED),
    ("•", "Refit the cross-modal map on augmented pairs → d2 0.749 → 0.63.", RED),
    ("•", "Augmentation on the real server → d2 → 0.56.", RED),
], size=12.5, gap=11)
box(s, 6.85, 3.0, 5.9, 3.3, fill=PANEL, round_=True)
text(s, 7.15, 3.2, 5.4, 0.4, [[("WHY (OUR READING)", 12, LIME, True)]])
bullets(s, 7.15, 3.7, 5.4, 2.5, [
    ("→", "Only ~350 base subjects — a couple of copies each multiplies samples, not subject diversity.", LIME),
    ("→", "Synthetic warps don't match real d2 transforms; the model overfits the synthetic distribution.", LIME),
    ("→", "Normalization beats augmentation: remove the distortion at inference, don't train to ignore it.", LIME),
], size=12, gap=9)

# ============================================================ 6. METRICS
s = slide()
header(s, "Results", LIME, "Honest, leak-free metrics")
rows = [
    ("Raw cosine, no Hungarian", "0.651", AMBER, 0.651),
    ("+ Hungarian assignment", "0.718", CYAN, 0.718),
    ("+ d3 grid feature (if legitimate)", "0.904", LIME, 0.904),
]
y = 2.05
for name, val, col, frac in rows:
    text(s, 0.65, y, 5.6, 0.4, [[(name, 14, WHITE, False)]], anchor=MSO_ANCHOR.MIDDLE)
    box(s, 6.3, y + 0.05, 5.3 * frac, 0.34, fill=col, round_=True)
    text(s, 11.75, y, 1.4, 0.4, [[(val, 15, col, True)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.62
text(s, 0.65, 4.0, 12.1, 0.4, [[("Per-dataset (leak-free)", 13, MUTE, True)]])
for i, (d, v, c) in enumerate([("dataset 1", "0.964", CYAN), ("dataset 2", "0.749  (wall)", MAG), ("dataset 3", "0.442  (bbox)", AMBER)]):
    xx = 0.65 + i * 4.15
    box(s, xx, 4.45, 3.9, 1.0, fill=PANEL, round_=True); box(s, xx, 4.45, 0.08, 1.0, fill=c)
    text(s, xx + 0.25, 4.58, 3.5, 0.4, [[(d, 13, MUTE, True)]])
    text(s, xx + 0.22, 4.9, 3.5, 0.5, [[(v, 22, c, True)]])
text(s, 0.65, 5.75, 12.1, 0.9, [[
    ("Hungarian is worth +0.068. The reproduced leak hits 0.931 — but our honest pipeline tops out at "
     "0.718 (strict) / 0.904 (d3-grid), and d2's elastic warp is the genuine remaining wall.",
     14, WHITE, False)]], line_spacing=1.2)

prs.save("EHL_Paris_Medical_Retrieval.pptx")
print("saved EHL_Paris_Medical_Retrieval.pptx", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
